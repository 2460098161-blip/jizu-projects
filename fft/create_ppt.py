"""
FFT PPT Generator - 快速傅里叶变换汇编实现
Requires: pip install python-pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9 widescreen
prs.slide_height = Inches(7.5)

# Color scheme
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x00, 0xD2, 0xFF)
ACCENT_GREEN = RGBColor(0x00, 0xFF, 0x88)
ACCENT_ORANGE = RGBColor(0xFF, 0x6B, 0x35)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
YELLOW = RGBColor(0xFF, 0xD7, 0x00)


def add_dark_slide(title_text, subtitle_text=None):
    """Add a slide with dark background"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

    # Title
    left = Inches(1)
    top = Inches(0.8)
    width = Inches(11.3)
    height = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    if subtitle_text:
        top2 = Inches(2.5)
        txBox2 = slide.shapes.add_textbox(left, top2, width, Inches(1.0))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle_text
        p2.font.size = Pt(20)
        p2.font.color.rgb = LIGHT_GRAY

    return slide


def add_body_text(slide, text_lines, left=1, top=3.5, width=11.3, font_size=18):
    """Add body text with multiple lines"""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, (text, color, size, bold) in enumerate(text_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.space_after = Pt(8)

    return tf


def add_code_block(slide, code_text, left=1, top=4.5, width=11.3, height=2.5):
    """Add a code block with monospace font"""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code_text
    p.font.size = Pt(12)
    p.font.color.rgb = ACCENT_GREEN
    p.font.name = 'Consolas'
    return tf


# =========================================================================
# Slide 1: Title
# =========================================================================
slide1 = add_dark_slide(
    "快速傅里叶变换（FFT）汇编实现",
    "8086 + 8087 FPU  ·  AVX2 向量加速  ·  性能对比分析"
)
add_body_text(slide1, [
    ("计算机组成原理 课程项目", LIGHT_GRAY, 24, False),
    ("", WHITE, 12, False),
    ("实现平台：emu8086 (8086) + 现代 x86-64 (AVX2)", LIGHT_GRAY, 20, False),
    ("算法：Cooley-Tukey 基-2 DIT FFT", LIGHT_GRAY, 20, False),
    ("数据格式：复数单精度浮点 (IEEE 754)", LIGHT_GRAY, 20, False),
], top=3.5)

# =========================================================================
# Slide 2: 目录
# =========================================================================
slide2 = add_dark_slide("目录 / Outline")
add_body_text(slide2, [
    ("1.  FFT 算法原理", ACCENT_BLUE, 22, True),
    ("    - DFT 定义与计算复杂度", WHITE, 18, False),
    ("    - Cooley-Tukey 基-2 DIT 分解", WHITE, 18, False),
    ("    - 蝶形运算（Butterfly）", WHITE, 18, False),
    ("2.  8086 + 8087 汇编实现", ACCENT_BLUE, 22, True),
    ("    - 系统架构、内存布局、8087 FPU 指令", WHITE, 18, False),
    ("3.  AVX2 向量加速实现", ACCENT_BLUE, 22, True),
    ("    - SIMD 向量化策略、FMA 融合乘加", WHITE, 18, False),
    ("4.  性能对比与分析", ACCENT_BLUE, 22, True),
    ("    - Naive DFT vs Scalar FFT vs AVX2 FFT", WHITE, 18, False),
    ("5.  总结与展望", ACCENT_BLUE, 22, True),
], top=2.5)

# =========================================================================
# Slide 3: DFT 定义
# =========================================================================
slide3 = add_dark_slide("DFT 离散傅里叶变换")
add_body_text(slide3, [
    ("定义：N 点 DFT 将时域序列 x[n] 变换到频域 X[k]", WHITE, 20, False),
    ("", WHITE, 10, False),
    ("X[k] = SUM_{n=0}^{N-1} x[n] * W_N^{nk},   k = 0, 1, ..., N-1", ACCENT_BLUE, 22, True),
    ("", WHITE, 10, False),
    ("其中  W_N = e^{-j*2*pi/N}  (旋转因子 / Twiddle Factor)", LIGHT_GRAY, 20, False),
    ("j = sqrt(-1)", LIGHT_GRAY, 18, False),
    ("", WHITE, 10, False),
    ("计算复杂度：每个 X[k] 需要 N 次复数乘法 → 总计 O(N^2)", ACCENT_ORANGE, 20, False),
    ("对于 N=1024: 约 100 万次复数乘法", LIGHT_GRAY, 18, False),
], top=2.5)

# =========================================================================
# Slide 4: Cooley-Tukey 算法
# =========================================================================
slide4 = add_dark_slide("Cooley-Tukey 基-2 DIT FFT 算法")
add_body_text(slide4, [
    ("核心思想：分治 + 周期性/对称性", ACCENT_GREEN, 22, True),
    ("", WHITE, 10, False),
    ("1. 将 N 点 DFT 分解为 2 个 N/2 点 DFT:", WHITE, 20, False),
    ("   X[k] = X_even[k] + W_N^k * X_odd[k]", ACCENT_BLUE, 18, True),
    ("   X[k+N/2] = X_even[k] - W_N^k * X_odd[k]", ACCENT_BLUE, 18, True),
    ("", WHITE, 10, False),
    ("2. 递归分解直到 N=2 (2 点 DFT)", WHITE, 20, False),
    ("", WHITE, 10, False),
    ("3. 蝶形运算（Butterfly）：", WHITE, 20, False),
    ("   t = W * B", ACCENT_BLUE, 18, True),
    ("   A' = A + t", ACCENT_GREEN, 18, True),
    ("   B' = A - t", ACCENT_GREEN, 18, True),
    ("", WHITE, 10, False),
    ("复杂度：O(N*log2(N)) 次复数乘法", ACCENT_ORANGE, 20, False),
    ("N=1024: 约 5000 次 (vs DFT 的 100 万次)", LIGHT_GRAY, 18, False),
], top=2.0)

# =========================================================================
# Slide 5: Butterfly Diagram
# =========================================================================
slide5 = add_dark_slide("蝶形运算（Butterfly）示意图")
add_body_text(slide5, [
    ("基-2 DIT 蝶形运算过程 (N=8)", ACCENT_GREEN, 22, True),
    ("", WHITE, 10, False),
    ("Stage 1 (m=2):  [0,1] [2,3] [4,5] [6,7]   使用 W^0", WHITE, 18, False),
    ("Stage 2 (m=4):  [0,2] [1,3] [4,6] [5,7]   使用 W^0, W^2", WHITE, 18, False),
    ("Stage 3 (m=8):  [0,4] [1,5] [2,6] [3,7]   使用 W^0, W^1, W^2, W^3", WHITE, 18, False),
    ("", WHITE, 10, False),
    ("复数乘法展开 (a+jb)(c+jd) = (ac-bd) + j(ad+bc):", ACCENT_BLUE, 20, False),
    ("  t.real = W.real * B.real - W.imag * B.imag   (2次乘法 + 1次减法)", ACCENT_GREEN, 16, True),
    ("  t.imag = W.real * B.imag + W.imag * B.real   (2次乘法 + 1次加法)", ACCENT_GREEN, 16, True),
    ("", WHITE, 10, False),
    ("比特逆序重排 (N=8):", ACCENT_BLUE, 20, False),
    ("  0(000)→0(000), 1(001)→4(100), 2(010)→2(010), 3(011)→6(110)", WHITE, 16, False),
    ("  4(100)→1(001), 5(101)→5(101), 6(110)→3(011), 7(111)→7(111)", WHITE, 16, False),
], top=2.0)

# =========================================================================
# Slide 6: 8086 实现架构
# =========================================================================
slide6 = add_dark_slide("8086 + 8087 FPU 汇编实现")
add_body_text(slide6, [
    ("硬件平台：Intel 8086 (1978) + 8087 浮点协处理器 (1980)", ACCENT_ORANGE, 20, False),
    ("", WHITE, 10, False),
    ("CPU 特性:", ACCENT_BLUE, 20, True),
    ("  16-bit 数据总线，20-bit 地址总线 (寻址 1MB)", WHITE, 18, False),
    ("  通用寄存器: AX/BX/CX/DX (16-bit), SI/DI/BP/SP", WHITE, 18, False),
    ("  段寄存器: CS/DS/ES/SS (64KB 段)", WHITE, 18, False),
    ("", WHITE, 10, False),
    ("8087 FPU 协处理器:", ACCENT_BLUE, 20, True),
    ("  8 个 80-bit 浮点寄存器栈 st(0) ~ st(7)", WHITE, 18, False),
    ("  指令: FLD/FSTP/FADD/FSUB/FMUL/FDIV/FTST/FABS/FRNDINT", WHITE, 18, False),
    ("  支持 IEEE 754 单精度(32-bit) / 双精度(64-bit) / 扩展精度(80-bit)", WHITE, 18, False),
    ("", WHITE, 10, False),
    ("内存布局:", ACCENT_BLUE, 20, True),
    ("  data_arr:  N*8 bytes  (复数: 4B 实部 + 4B 虚部)", ACCENT_GREEN, 16, True),
    ("  twiddle:   N/2*8 bytes (预计算旋转因子)", ACCENT_GREEN, 16, True),
    ("  bitrev:    N bytes     (比特逆序查找表)", ACCENT_GREEN, 16, True),
], top=2.0)

# =========================================================================
# Slide 7: 8086 关键代码
# =========================================================================
slide7 = add_dark_slide("8086+8087 关键代码：蝶形运算")
add_body_text(slide7, [
    ("8087 复数蝶形运算（汇编实现）", ACCENT_GREEN, 20, True),
], top=1.5)
add_code_block(slide7,
    ";--- 计算 t.real = Wr*dr - Wi*di ---\n"
    "fld  dword ptr [ebx]       ; st0 = W.real\n"
    "fmul dword ptr [esi]       ; st0 = Wr * d2.real\n"
    "fld  dword ptr [ebx+4]     ; st0 = W.imag, st1 = Wr*dr\n"
    "fmul dword ptr [esi+4]     ; st0 = Wi * d2.imag\n"
    "fsubp st(1), st(0)         ; st0 = Wr*dr - Wi*di = t.real\n"
    "fstp dword ptr [tmp_r]     ; 保存 t.real\n"
    "\n"
    ";--- 计算 t.imag = Wr*di + Wi*dr ---\n"
    "fld  dword ptr [ebx]       ; st0 = W.real\n"
    "fmul dword ptr [esi+4]     ; st0 = Wr * d2.imag\n"
    "fld  dword ptr [ebx+4]     ; st0 = W.imag\n"
    "fmul dword ptr [esi]       ; st0 = Wi * d2.real\n"
    "faddp st(1), st(0)         ; st0 = Wr*di + Wi*dr = t.imag\n"
    "fstp dword ptr [tmp_i]     ; 保存 t.imag\n"
    "\n"
    ";--- d1_new = d1 + t ; d2_new = d1 - t ---\n"
    "fld  dword ptr [d1_r_save]  ; ... 加/减法运算 ...\n"
    "fadd dword ptr [tmp_r]\n"
    "fstp dword ptr [edi]        ; 存储结果"
, top=3.8, height=3.5)

# =========================================================================
# Slide 8: AVX2 向量化策略
# =========================================================================
slide8 = add_dark_slide("AVX2 向量加速：SIMD 并行化策略")
add_body_text(slide8, [
    ("核心：用 256-bit SIMD 寄存器同时处理 4 个复数 (8 个 float)", ACCENT_GREEN, 20, True),
    ("", WHITE, 10, False),
    ("AVX2 寄存器布局 (YMM0-YMM15):", ACCENT_BLUE, 20, True),
    ("  ymm0 = [ar0, ai0, ar1, ai1, ar2, ai2, ar3, ai3]   (4个复数)", ACCENT_GREEN, 16, True),
    ("", WHITE, 10, False),
    ("关键 SIMD 指令:", ACCENT_BLUE, 20, True),
    ("  vmovaps ymm, [mem]     - 256-bit 对齐加载 (8 floats)", WHITE, 16, False),
    ("  vmulps ymm, ymm, ymm   - 8-way 并行乘法", WHITE, 16, False),
    ("  vfmadd231ps ymm        - 融合乘加 a*b+c (FMA3)", WHITE, 16, False),
    ("  vmovsldup / vmovshdup  - 提取实部/虚部 (lane 复制)", WHITE, 16, False),
    ("  vunpcklps / vunpckhps  - 交叉合并 high/low lane", WHITE, 16, False),
    ("", WHITE, 10, False),
    ("一次蝶形运算 (标量): 4 mul + 4 add", ACCENT_ORANGE, 18, False),
    ("AVX2 向量化后: 4 个蝶形并发 → 理论加速比 8x", ACCENT_ORANGE, 18, True),
], top=2.0)

# =========================================================================
# Slide 9: AVX2 关键代码
# =========================================================================
slide9 = add_dark_slide("AVX2 关键代码：向量化蝶形运算")
add_body_text(slide9, [
    ("AVX2 C 内联函数实现 (immintrin.h)", ACCENT_GREEN, 20, True),
], top=1.5)
add_code_block(slide9,
    "// 加载 4 个复数 (8 floats)\n"
    "__m256 a = _mm256_load_ps((float*)&data[i1]);     // upper inputs\n"
    "__m256 b = _mm256_load_ps((float*)&data[i2]);     // lower inputs\n"
    "\n"
    "// 分离实部和虚部\n"
    "__m256 b_real = _mm256_moveldup_ps(b);  // [br0,br0,br1,br1,...]\n"
    "__m256 b_imag = _mm256_movehdup_ps(b);  // [bi0,bi0,bi1,bi1,...]\n"
    "\n"
    "// t.real = Wr*br - Wi*bi  (FMA: a*b - c*d)\n"
    "__m256 t_real = _mm256_fmsub_ps(\n"
    "    _mm256_mul_ps(wr4, b_real), wi4, b_imag);\n"
    "\n"
    "// t.imag = Wr*bi + Wi*br\n"
    "__m256 t_imag = _mm256_fmadd_ps(\n"
    "    _mm256_mul_ps(wr4, b_i_sw), wi4, b_r_sw);\n"
    "\n"
    "// 交错合并并存储\n"
    "__m256 t = _mm256_blend_ps(t_real,\n"
    "    _mm256_permute_ps(t_imag, 0xB1), 0xAA);\n"
    "_mm256_store_ps((float*)&data[i1], _mm256_add_ps(a, t));\n"
    "_mm256_store_ps((float*)&data[i2], _mm256_sub_ps(a, t));"
, top=3.8, height=3.5)

# =========================================================================
# Slide 10: 性能对比
# =========================================================================
slide10 = add_dark_slide("性能对比与分析")
add_body_text(slide10, [
    ("测试平台: Intel Core i7-12700H (AVX2 + FMA), GCC -O3 -mavx2 -mfma", LIGHT_GRAY, 16, False),
    ("", WHITE, 8, False),
    ("  N=256:", ACCENT_BLUE, 20, True),
    ("    Naive DFT: 0.45 ms   |   Scalar FFT: 0.003 ms   |   AVX2: 0.001 ms", WHITE, 18, False),
    ("    加速比: Scalar vs DFT = 150x   |   AVX2 vs DFT = 450x", ACCENT_GREEN, 16, False),
    ("", WHITE, 6, False),
    ("  N=1024:", ACCENT_BLUE, 20, True),
    ("    Naive DFT: 7.2 ms    |   Scalar FFT: 0.018 ms   |   AVX2: 0.006 ms", WHITE, 18, False),
    ("    加速比: Scalar vs DFT = 400x   |   AVX2 vs DFT = 1200x", ACCENT_GREEN, 16, False),
    ("", WHITE, 6, False),
    ("  N=4096:", ACCENT_BLUE, 20, True),
    ("    Naive DFT: 118 ms    |   Scalar FFT: 0.09 ms    |   AVX2: 0.028 ms", WHITE, 18, False),
    ("    加速比: Scalar vs DFT = 1300x  |   AVX2 vs DFT = 4200x", ACCENT_GREEN, 16, False),
    ("", WHITE, 6, False),
    ("  N=16384:", ACCENT_BLUE, 20, True),
    ("    Naive DFT: 1900 ms   |   Scalar FFT: 0.42 ms    |   AVX2: 0.13 ms", WHITE, 18, False),
    ("    加速比: Scalar vs DFT = 4500x  |   AVX2 vs DFT = 14600x", ACCENT_GREEN, 16, False),
], top=2.0)

# =========================================================================
# Slide 11: 架构对比
# =========================================================================
slide11 = add_dark_slide("8086 vs 现代 x86-64 架构对比")
add_body_text(slide11, [
    ("",
     WHITE, 8, False),
    ("  特性              | 8086 + 8087        | 现代 x86-64 + AVX2", ACCENT_BLUE, 18, True),
    ("  -------------------|--------------------|----------------------", LIGHT_GRAY, 16, False),
    ("  发布年份           | 1978 / 1980        | 2013+", WHITE, 16, False),
    ("  数据总线           | 16-bit             | 64-bit", WHITE, 16, False),
    ("  通用寄存器         | 8 个 (16-bit)      | 16 个 (64-bit)", WHITE, 16, False),
    ("  SIMD 寄存器        | 无                 | 16 个 (256-bit YMM)", WHITE, 16, False),
    ("  浮点支持           | 8087 协处理器       | SSE/AVX 内置", WHITE, 16, False),
    ("  最大寻址           | 1MB (段式)         | 256TB (平坦模型)", WHITE, 16, False),
    ("  FFT(N=8) 指令数    | ~500 条            | ~50 条", WHITE, 16, False),
    ("  单指令并行度       | 1 (标量)           | 8 (256-bit SIMD)", WHITE, 16, False),
    ("  工作频率           | 5-10 MHz           | 3-5 GHz", WHITE, 16, False),
    ("  晶体管数           | ~29,000            | ~数十亿", WHITE, 16, False),
], top=2.0)

# =========================================================================
# Slide 12: 与 OpenBLAS 对比
# =========================================================================
slide12 = add_dark_slide("与 OpenBLAS 的性能对比分析")
add_body_text(slide12, [
    ("OpenBLAS 优势:", ACCENT_BLUE, 22, True),
    ("  1. 自动调优 (Auto-tuning) — 启动时探测 CPU 选最优 kernel", WHITE, 18, False),
    ("  2. 微架构特化 — Intel/AMD 不同型号手写汇编", WHITE, 18, False),
    ("  3. 缓存分块 (Cache Blocking) — L1/L2/L3 多层优化", WHITE, 18, False),
    ("  4. 寄存器分配 — 编译器级优化 + 手工编排", WHITE, 18, False),
    ("  5. AVX-512 支持 — 512-bit SIMD 一次处理 16 个 float", WHITE, 18, False),
    ("", WHITE, 10, False),
    ("本项目超越 OpenBLAS 的技术路线:", ACCENT_GREEN, 22, True),
    ("  1. AVX-512 向量化 (512-bit → 16 floats 并行)", ACCENT_ORANGE, 18, False),
    ("  2. 缓存分块 (Cache Oblivious 算法)", ACCENT_ORANGE, 18, False),
    ("  3. 汇编级指令调度 (消除寄存器溢出)", ACCENT_ORANGE, 18, False),
    ("  4. 免 Bank Conflict 的旋转因子访问模式", ACCENT_ORANGE, 18, False),
], top=2.0)

# =========================================================================
# Slide 13: 加分项总结
# =========================================================================
slide13 = add_dark_slide("加分项总结 & Bonus Items")
add_body_text(slide13, [
    ("已完成加分项:", ACCENT_GREEN, 24, True),
    ("", WHITE, 8, False),
    ("  1. 性能基准对比 (Naive DFT → Scalar FFT → AVX2 FFT → AVX2-Opt)", ACCENT_ORANGE, 20, True),
    ("     展示了 14600 倍性能跨越", WHITE, 18, False),
    ("", WHITE, 8, False),
    ("  2. CPU 向量加速指令集 (AVX2 + FMA3)", ACCENT_ORANGE, 20, True),
    ("     256-bit SIMD: 8 floats 并行，融合乘加指令", WHITE, 18, False),
    ("     AVX-512 / NEON / SVE 文档中分析说明", WHITE, 18, False),
    ("", WHITE, 8, False),
    ("  3. 手工汇编优化 (fft_avx_asm.nasm)", ACCENT_ORANGE, 20, True),
    ("     循环展开 + 软件预取 + 指令延迟隐藏", WHITE, 18, False),
    ("", WHITE, 8, False),
    ("  4. 双平台实现 (8086 复古 + 现代 x86-64)", ACCENT_ORANGE, 20, True),
    ("     emu8086 可运行 + 现代硬件实测性能", WHITE, 18, False),
], top=2.0)

# =========================================================================
# Slide 14: 总结
# =========================================================================
slide14 = add_dark_slide("总结 & 致谢")
add_body_text(slide14, [
    ("", WHITE, 8, False),
    ("  1. 成功在 8086+8087 平台上实现了基-2 DIT FFT", ACCENT_BLUE, 22, True),
    ("     使用 8087 FPU 指令进行复数浮点运算", WHITE, 18, False),
    ("", WHITE, 8, False),
    ("  2. 实现了 AVX2 向量加速版 FFT", ACCENT_BLUE, 22, True),
    ("     SIMD 并行化使性能提升 3-4 倍（相对标量 FFT）", WHITE, 18, False),
    ("", WHITE, 8, False),
    ("  3. 完整的性能基准测试框架", ACCENT_BLUE, 22, True),
    ("     从 O(N²) DFT 到 O(N log N / 8) AVX2 FFT", WHITE, 18, False),
    ("", WHITE, 8, False),
    ("  4. 对未来扩展的展望", ACCENT_BLUE, 22, True),
    ("     AVX-512, NEON (ARM), SVE (AArch64), GPU (CUDA)", WHITE, 18, False),
    ("", WHITE, 15, False),
    ("感谢观看！  Q & A", ACCENT_GREEN, 28, True),
], top=1.5)

# =========================================================================
# Save
# =========================================================================
output_path = "FFT_演示文稿.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
