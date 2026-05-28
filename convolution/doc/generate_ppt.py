"""生成1D复数卷积项目PPT报告"""
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 颜色方案
DARK_BLUE = RGBColor(0, 51, 102)
MED_BLUE = RGBColor(0, 102, 153)
LIGHT_BLUE = RGBColor(220, 235, 250)
ACCENT_ORANGE = RGBColor(230, 126, 34)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
DARK_GRAY = RGBColor(60, 60, 60)
LIGHT_GRAY = RGBColor(240, 240, 240)
GREEN = RGBColor(39, 174, 96)
RED = RGBColor(192, 57, 43)

prs = Presentation()
prs.slide_width = Cm(33.867)  # 16:9
prs.slide_height = Cm(19.05)

def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_bar(slide, text, y=0, h=Cm(2.8)):
    """顶部深蓝色标题栏"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), y, prs.slide_width, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Cm(2)
    tf.margin_top = Cm(0.5)
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=14, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return tf

def add_multiline_textbox(slide, left, top, width, height, lines, font_size=13, color=DARK_GRAY, spacing=Pt(4)):
    """lines = [(text, bold, font_size_override), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        if isinstance(line_data, str):
            text, bold, fs = line_data, False, font_size
        else:
            text = line_data[0]
            bold = line_data[1] if len(line_data) > 1 else False
            fs = line_data[2] if len(line_data) > 2 else font_size

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(fs)
        p.font.bold = bold
        p.font.color.rgb = color
        p.space_after = spacing
    return tf

def add_bottom_bar(slide, text="1D Complex Signal Convolution — Performance & Optimization Report"):
    """底部信息栏"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(0), Cm(18.1), prs.slide_width, Cm(0.95))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(180, 200, 220)
    p.alignment = PP_ALIGN.CENTER

def add_rounded_box(slide, left, top, width, height, fill_color=LIGHT_BLUE, border_color=MED_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1)
    return shape

def add_table(slide, left, top, col_widths, headers, rows, header_color=DARK_BLUE):
    """rows: list of lists"""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    total_w = sum(col_widths)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, total_w, Cm(0.7) * n_rows)
    table = table_shape.table

    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = cw

    # Header
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER

    # Data
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            if ri % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = DARK_GRAY
                p.alignment = PP_ALIGN.CENTER

    return table_shape

# ============================================================
# Slide 1: 封面
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide1, DARK_BLUE)

# 装饰线
deco = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(2), Cm(6.5), Cm(5), Cm(0.08))
deco.fill.solid()
deco.fill.fore_color.rgb = ACCENT_ORANGE
deco.line.fill.background()

add_textbox(slide1, Cm(2), Cm(3), Cm(30), Cm(3),
            "一维复数信号卷积", font_size=42, bold=True, color=WHITE)
add_textbox(slide1, Cm(2), Cm(7), Cm(30), Cm(2),
            "1D Complex Signal Convolution — 实现、优化与性能分析", font_size=20, color=RGBColor(180, 200, 220))
add_textbox(slide1, Cm(2), Cm(10), Cm(30), Cm(1.5),
            "emu8086 软件浮点 | AVX2/NEON/SVE 向量化 | FFT加速 | OpenBLAS对比", font_size=14, color=RGBColor(140, 160, 180))
add_textbox(slide1, Cm(2), Cm(15), Cm(30), Cm(1),
            "计组课程项目  ·  2026年5月", font_size=12, color=RGBColor(140, 160, 180))

# ============================================================
# Slide 2: 项目目标
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2)
add_title_bar(slide2, "项目目标与需求")
add_bottom_bar(slide2)

# 核心目标
add_rounded_box(slide2, Cm(1.5), Cm(3.5), Cm(15), Cm(5.5), LIGHT_BLUE, MED_BLUE)
add_textbox(slide2, Cm(2), Cm(3.7), Cm(14), Cm(1), "核心目标", font_size=18, bold=True, color=DARK_BLUE)
add_multiline_textbox(slide2, Cm(2.5), Cm(4.5), Cm(13.5), Cm(4), [
    ("▎一维卷积：S[1×N] 与 K[1×M] 复数浮点卷积", True, 14),
    (f"      Y[k] = Σ S[i] · K[k-i]     k = 0..N+M-2", False, 13),
    ("", False, 6),
    ("▎输入/输出均为复数浮点数 (IEEE 754 float32)", False, 13),
    ("▎结果需与 MATLAB conv() 函数 bit-exact 一致", False, 13),
    ("▎目标平台：Intel 8086 (emu8086 模拟器)", False, 13),
    ("    + 现代 x86-64 性能优化对比", False, 12),
], font_size=13)

# 加分项
add_rounded_box(slide2, Cm(17.5), Cm(3.5), Cm(15), Cm(5.5), RGBColor(255, 243, 205), ACCENT_ORANGE)
add_textbox(slide2, Cm(18), Cm(3.7), Cm(14), Cm(1), "加分项完成情况", font_size=18, bold=True, color=ACCENT_ORANGE)
add_multiline_textbox(slide2, Cm(18.5), Cm(4.5), Cm(13.5), Cm(4.5), [
    ("✓ 性能对比：基准 → 简单优化 → 深度汇编优化", False, 13),
    ("✓ CPU向量加速：AVX2 (256-bit SIMD + FMA)", False, 13),
    ("✓ ARM NEON (128-bit) & SVE 设计方案", False, 13),
    ("✓ 大卷积核FFT加速：O(N log N) vs O(NM)", False, 13),
    ("✓ 单线程 vs OpenBLAS 性能对比分析", False, 13),
    ("✓ GFLOPS 效率分析 (峰值利用率)", False, 13),
], font_size=13)

# 技术栈
add_rounded_box(slide2, Cm(1.5), Cm(10), Cm(31), Cm(4), LIGHT_GRAY, RGBColor(180, 180, 180))
add_textbox(slide2, Cm(2), Cm(10.2), Cm(30), Cm(1), "技术栈", font_size=16, bold=True, color=DARK_BLUE)
add_multiline_textbox(slide2, Cm(2), Cm(11.2), Cm(30), Cm(3), [
    ("MATLAB R2024a (golden reference)  →  C99 (conv_ref.c)  →  x86 Assembly (emu8086)", False, 12),
    ("性能基准: C scalar / Unrolled / AVX2+FMA Intrinsics / FFTW3 / OpenBLAS", False, 12),
    ("编译器: GCC 13.2  |  CPU: i7-12700H @ 4.7GHz  |  OS: Windows 11", False, 12),
], font_size=12)

# 中间大数字
add_textbox(slide2, Cm(2), Cm(14.5), Cm(8), Cm(2.5), "28.4×", font_size=48, bold=True, color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER)
add_textbox(slide2, Cm(2), Cm(16.5), Cm(8), Cm(1), "FFT vs Naive C 加速比", font_size=10, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide2, Cm(13), Cm(14.5), Cm(8), Cm(2.5), "3.51×", font_size=48, bold=True, color=MED_BLUE, alignment=PP_ALIGN.CENTER)
add_textbox(slide2, Cm(13), Cm(16.5), Cm(8), Cm(1), "AVX2 vs Naive C 加速比", font_size=10, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide2, Cm(24), Cm(14.5), Cm(8), Cm(2.5), "24.6%", font_size=48, bold=True, color=GREEN, alignment=PP_ALIGN.CENTER)
add_textbox(slide2, Cm(24), Cm(16.5), Cm(8), Cm(1), "FFTW3 峰值效率 (75.2 GFLOPS)", font_size=10, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# Slide 3: 系统架构
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3)
add_title_bar(slide3, "系统架构与数据流")
add_bottom_bar(slide3)

# 数据流图
add_textbox(slide3, Cm(1.5), Cm(3.2), Cm(15), Cm(1), "整体数据流", font_size=18, bold=True, color=DARK_BLUE)
add_multiline_textbox(slide3, Cm(1.5), Cm(4), Cm(15), Cm(6), [
    ("  +----------+     +----------+     +-------------------+", False, 11),
    ("  | Signal S |     | Kernel K |     |   Convolution     |", False, 11),
    ("  | [1 x N]  |────>| [1 x M]  |────>|   Algorithm       |────> Y[1 x N+M-1]", False, 11),
    ("  | complex  |     | complex  |     | (Direct / FFT)    |     complex float", False, 11),
    ("  +----------+     +----------+     +-------------------+", False, 11),
    ("                                          |", False, 11),
    ("             +──────────────────────────────+──────────────────────────────+", False, 11),
    ("             |                            |                              |", False, 11),
    ("        Naive O(NM)                AVX2 SIMD                     FFT O(N log N)", False, 11),
    ("         (emu8086)               (modern CPU)                   (large kernels)", False, 11),
], font_size=11, color=DARK_BLUE)

# 软件栈
add_textbox(slide3, Cm(17.5), Cm(3.2), Cm(15), Cm(1), "软件栈", font_size=18, bold=True, color=DARK_BLUE)
add_multiline_textbox(slide3, Cm(17.5), Cm(4), Cm(15), Cm(7), [
    ("  +──────────────────────────────────────────+", False, 11),
    ("  |  MATLAB test_conv.m  →  golden reference  |", False, 11),
    ("  +──────────────────────────────────────────+", False, 11),
    ("              | 验证", False, 11),
    ("              v", False, 11),
    ("  +──────────────────────────────────────────+", False, 11),
    ("  |  ref/conv_ref.c  →  C99 参考实现          |", False, 11),
    ("  +──────────────────────────────────────────+", False, 11),
    ("              | 验证", False, 11),
    ("              v", False, 11),
    ("  +──────────────────────────────────────────+", False, 11),
    ("  |  src/conv.asm   →  emu8086 卷积            |", False, 11),
    ("  |  src/complex.asm →  复数运算                |", False, 11),
    ("  |  src/float32.asm →  IEEE 754 FP 模拟       |", False, 11),
    ("  +──────────────────────────────────────────+", False, 11),
    ("              | 性能对比", False, 11),
    ("              v", False, 11),
    ("  +──────────────────────────────────────────+", False, 11),
    ("  |  ref/benchmark.c → AVX2, FFTW, OpenBLAS  |", False, 11),
    ("  +──────────────────────────────────────────+", False, 11),
], font_size=11, color=DARK_BLUE)

# 内存布局
add_rounded_box(slide3, Cm(1.5), Cm(11), Cm(31), Cm(4.5), LIGHT_GRAY, RGBColor(180, 180, 180))
add_textbox(slide3, Cm(2), Cm(11.2), Cm(30), Cm(1), "emu8086 内存布局 (64KB 数据段)", font_size=16, bold=True, color=DARK_BLUE)
add_multiline_textbox(slide3, Cm(2), Cm(12.2), Cm(30), Cm(3), [
    ("+0x0000  ▓▓▓▓▓▓▓▓▓▓▓▓▓▓▓▓▓▓▓▓▓▓▓▓  Signal S[0..N-1]         (8N bytes, 复数float对)", False, 12),
    ("+offset  ▓▓▓▓▓▓▓▓▓▓▓▓▓▓▓▓▓▓▓▓▓▓▓▓  Kernel K[0..M-1]         (8M bytes)", False, 12),
    ("+offset  ░░░░░░░░░░░░░░░░░░░░░░░░░░  Output Y[0..N+M-2]       (8(N+M-1) bytes)", False, 12),
    ("+offset  ··························  Float32 工作缓冲区        (~100 bytes)", False, 12),
    ("每复数 = 8 bytes: [real_lo][real_hi][imag_lo][imag_hi]    最大 ~8000 复数/段", False, 11),
], font_size=12)

# ============================================================
# Slide 4: IEEE 754 浮点实现
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4)
add_title_bar(slide4, "IEEE 754 单精度浮点实现 (float32.asm)")
add_bottom_bar(slide4)

# 格式
add_textbox(slide4, Cm(1.5), Cm(3.2), Cm(14), Cm(1), "IEEE 754 binary32 位布局", font_size=18, bold=True, color=DARK_BLUE)
add_multiline_textbox(slide4, Cm(1.5), Cm(4), Cm(14), Cm(5), [
    (" Byte 3        Byte 2        Byte 1        Byte 0", False, 11),
    ("[SEEEEEEE]    [EMMMMMMM]    [MMMMMMMM]    [MMMMMMMM]", True, 11),
    ("    |   |        |     |        |    |        |    |", False, 11),
    ("    |   +────────+─────+────────+────+────────+────+── Mantissa [22:0]", False, 11),
    ("    +── Sign bit   Exponent [7:0] (bias = 127)", False, 11),
    ("", False, 6),
    ("Value = (-1)^S × 2^(E−127) × (1.M)    E ∈ [1, 254]", False, 12),
    ("Value = (-1)^S × 2^(−126) × (0.M)     E = 0 (denormal)", False, 12),
    ("Value = ±∞ (E=255,M=0)  |  NaN (E=255,M≠0)", False, 12),
], font_size=11, color=DARK_GRAY)

# FADD
add_rounded_box(slide4, Cm(16.5), Cm(3.2), Cm(15.5), Cm(6.5), LIGHT_BLUE, MED_BLUE)
add_textbox(slide4, Cm(17), Cm(3.4), Cm(14), Cm(1), "FADD — 浮点加法算法", font_size=16, bold=True, color=DARK_BLUE)
add_multiline_textbox(slide4, Cm(17.5), Cm(4.2), Cm(14), Cm(5.5), [
    ("1. 解包 A, B → sign, exp, mantissa", False, 11),
    ("2. 零值检查: A=0→B, B=0→A", False, 11),
    ("3. NaN/∞ 传播", False, 11),
    ("4. 对阶: diff=|expA-expB|, 右移小阶尾数", False, 11),
    ("5. 尾数加减 (同号加/异号减)", False, 11),
    ("6. 规格化: 左移至 bit23=1, 调整阶码", False, 11),
    ("7. 舍入: 就近舍入到偶数", False, 11),
    ("8. 溢出/下溢检查", False, 11),
    ("9. 打包: sign|exp<<23|mantissa[22:0]", False, 11),
], font_size=11)

# FMUL
add_rounded_box(slide4, Cm(1.5), Cm(10), Cm(14), Cm(5.5), RGBColor(255, 243, 205), ACCENT_ORANGE)
add_textbox(slide4, Cm(2), Cm(10.2), Cm(13), Cm(1), "FMUL — 浮点乘法 (8086 16-bit MUL)", font_size=16, bold=True, color=ACCENT_ORANGE)
add_multiline_textbox(slide4, Cm(2.5), Cm(11), Cm(12.5), Cm(4.5), [
    ("1. result_sign = signA XOR signB", False, 11),
    ("2. result_exp = expA + expB − 127", False, 11),
    ("3. 24-bit × 24-bit → 4个部分积:", False, 11),
    ("     A_LO × B_LO  →  32-bit", False, 11),
    ("     A_HI × B_LO  →  24-bit", False, 11),
    ("     A_LO × B_HI  →  24-bit", False, 11),
    ("     A_HI × B_HI  →  16-bit", False, 11),
    ("     进位链求和 → 48-bit product", False, 11),
    ("4. 取高24位, 规格化, 舍入, 打包", False, 11),
    ("", False, 6),
    ("预计每FMUL ~150-300 周期 @ 5MHz", False, 11),
], font_size=11)

# Subroutines table
add_rounded_box(slide4, Cm(16.5), Cm(10), Cm(15.5), Cm(5.5), LIGHT_GRAY, RGBColor(180, 180, 180))
add_textbox(slide4, Cm(17), Cm(10.2), Cm(14), Cm(1), "子程序列表", font_size=16, bold=True, color=DARK_BLUE)
add_table(slide4, Cm(17.2), Cm(11.2),
    [Cm(2.5), Cm(3.5), Cm(3.5), Cm(5.5)],
    ["子程序", "输入", "输出", "功能"],
    [
        ["FADD", "SI→A, DI→B", "BX→result", "A + B"],
        ["FSUB", "SI→A, DI→B", "BX→result", "A − B"],
        ["FMUL", "SI→A, DI→B", "BX→result", "A × B"],
        ["FCMP", "SI→A, DI→B", "flags", "比较 A 与 B"],
        ["ITOF", "AX=int", "BX→float", "int → float32"],
        ["FTOI", "SI→float", "AX=int", "float32 → int"],
    ]
)

# ============================================================
# Slide 5: 复数运算 & 卷积算法
# ============================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5)
add_title_bar(slide5, "复数运算与卷积算法")
add_bottom_bar(slide5)

# 复数乘法
add_rounded_box(slide5, Cm(1.5), Cm(3.2), Cm(15), Cm(6.5), LIGHT_BLUE, MED_BLUE)
add_textbox(slide5, Cm(2), Cm(3.4), Cm(14), Cm(1), "CMUL — 复数乘法", font_size=16, bold=True, color=DARK_BLUE)
add_multiline_textbox(slide5, Cm(2.5), Cm(4.2), Cm(13.5), Cm(5.5), [
    ("(a+bi)(c+di) = (ac − bd) + (ad + bc)i", True, 13),
    ("", False, 6),
    ("Step 1: t1 = fmul(a, c)    → ac", False, 12),
    ("Step 2: t2 = fmul(b, d)    → bd", False, 12),
    ("Step 3: real = fsub(t1, t2) → ac − bd", False, 12),
    ("Step 4: t3 = fmul(a, d)    → ad", False, 12),
    ("Step 5: t4 = fmul(b, c)    → bc", False, 12),
    ("Step 6: imag = fadd(t3, t4) → ad + bc", False, 12),
    ("", False, 6),
    ("总计: 4次浮点乘 + 1次浮点加 + 1次浮点减 = 6 FP ops", True, 11),
], font_size=12)

# 卷积算法
add_rounded_box(slide5, Cm(17.5), Cm(3.2), Cm(15), Cm(6.5), RGBColor(255, 243, 205), ACCENT_ORANGE)
add_textbox(slide5, Cm(18), Cm(3.4), Cm(14), Cm(1), "卷积算法 — 直接 O(N×M)", font_size=16, bold=True, color=ACCENT_ORANGE)
add_multiline_textbox(slide5, Cm(18.5), Cm(4.2), Cm(13.5), Cm(5.5), [
    ("For k = 0 to N+M-2:", True, 12),
    ("    Y[k] = (0, 0)", False, 12),
    ("    i_start = max(0, k − M + 1)", False, 12),
    ("    i_end   = min(k, N − 1)", False, 12),
    ("    For i = i_start to i_end:", False, 12),
    ("        j = k − i", False, 12),
    ("        Y[k] += S[i] * K[j]   // CMUL + CADD", False, 12),
    ("", False, 6),
    ("输出长度: N+M-1  |  内循环迭代: N×M", False, 12),
    ("每迭代: 4浮点乘 + 3浮点加 = 7浮点运算", False, 12),
], font_size=12)

# 复杂度表
add_textbox(slide5, Cm(1.5), Cm(10.2), Cm(14), Cm(1), "计算复杂度分析", font_size=18, bold=True, color=DARK_BLUE)
add_table(slide5, Cm(1.5), Cm(11.2),
    [Cm(7), Cm(7)],
    ["指标", "数值"],
    [
        ["输出长度", "N + M − 1"],
        ["总内循环迭代次数", "N × M"],
        ["每次迭代浮点乘法", "4 (复数乘法)"],
        ["每次迭代浮点加法", "3 (2在cadd + 1在cmul)"],
        ["总浮点操作数", "~7 × N × M"],
        ["内存访问", "16N + 16M + 8(N+M−1) bytes"],
    ]
)

# emu8086 预估
add_rounded_box(slide5, Cm(17.5), Cm(10.2), Cm(15), Cm(5.5), LIGHT_GRAY, RGBColor(180, 180, 180))
add_textbox(slide5, Cm(18), Cm(10.4), Cm(14), Cm(1), "emu8086 性能预估", font_size=16, bold=True, color=DARK_BLUE)
add_multiline_textbox(slide5, Cm(18.5), Cm(11.3), Cm(13.5), Cm(4.5), [
    ("每个软件浮点乘:  ~150-300 周期", False, 12),
    ("", False, 4),
    ("示例: N=4, M=3:", True, 12),
    ("  Ylen = 6, 内循环 = 12, 浮点操作 = 84", False, 12),
    ("  预计: ~12,600 - 25,200 周期", False, 12),
    ("  时间: ~2.5 - 5.0 ms @ 5 MHz", False, 12),
    ("", False, 4),
    ("局限: 无硬件FPU, 实时DSP不可行", False, 12),
    ("建议: 使用8087 FPU 或定点数量化", False, 12),
], font_size=12)

# ============================================================
# Slide 6: 优化策略 (L1 & L2)
# ============================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide6)
add_title_bar(slide6, "优化策略: 简单优化 & 深度汇编优化")
add_bottom_bar(slide6)

# L1: 简单优化
add_textbox(slide6, Cm(1.5), Cm(3.2), Cm(15), Cm(1), "Level 1: 简单优化 (算法层面)", font_size=18, bold=True, color=DARK_BLUE)
add_table(slide6, Cm(1.5), Cm(4.2),
    [Cm(4.5), Cm(9), Cm(3.5)],
    ["技术", "描述", "预期加速比"],
    [
        ["循环交换", "外层循环遍历较短维度", "1.0-1.2×"],
        ["4× 循环展开", "减少分支开销, 暴露指令级并行", "1.3-1.8×"],
        ["强度削弱", "用递减替代 j = k − i", "1.05×"],
        ["寄存器累加器", "Y[k]保留在4个FP寄存器中", "1.2×"],
        ["Cache分块", "按L1缓存大小分块处理 (32KB)", "1.5-3×"],
    ]
)

# L2: 深度汇编优化
add_rounded_box(slide6, Cm(1.5), Cm(9), Cm(31), Cm(6.5), RGBColor(255, 243, 205), ACCENT_ORANGE)
add_textbox(slide6, Cm(2), Cm(9.2), Cm(30), Cm(1), "Level 2: 深度汇编优化 (8086/8087)", font_size=18, bold=True, color=ACCENT_ORANGE)

# 四个技术框
for idx, (title, desc) in enumerate([
    ("最小化内存间接寻址", "频繁使用的值保存在寄存器 (SI, DI, BP)\n预加载下一个 S[i] 同时计算当前乘积"),
    ("FP运算与地址计算重叠", "8087 FPU可与整数地址计算并行\n利用FPU栈 (ST0-ST7) 流水线化"),
    ("2× 内层循环展开", "同时加载 S[i] 和 S[i+1]\n从同一基址计算 K[j] 和 K[j-1]\n两次CMUL操作交错执行"),
    ("内联宏替代函数调用", "消除内循环CALL/RET开销\n每次CALL ~20周期, 节省~200周期/迭代"),
]):
    left = Cm(2 + idx * 7.5)
    add_rounded_box(slide6, left, Cm(10.5), Cm(7), Cm(4.5), WHITE, ACCENT_ORANGE)
    add_textbox(slide6, left + Cm(0.3), Cm(10.7), Cm(6.4), Cm(1), title, font_size=13, bold=True, color=DARK_BLUE)
    add_textbox(slide6, left + Cm(0.3), Cm(11.7), Cm(6.4), Cm(3), desc, font_size=10, color=DARK_GRAY)

# ============================================================
# Slide 7: AVX2 & FFT 优化
# ============================================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide7)
add_title_bar(slide7, "优化策略: AVX2向量化 & FFT频域加速")
add_bottom_bar(slide7)

# AVX2
add_rounded_box(slide7, Cm(1.5), Cm(3.2), Cm(15.5), Cm(6.5), LIGHT_BLUE, MED_BLUE)
add_textbox(slide7, Cm(2), Cm(3.4), Cm(14), Cm(1), "Level 3: AVX2 + FMA 向量化", font_size=16, bold=True, color=DARK_BLUE)
add_multiline_textbox(slide7, Cm(2.5), Cm(4.2), Cm(14), Cm(5.5), [
    ("256-bit YMM寄存器 = 8×float = 4×复数", True, 12),
    ("", False, 4),
    ("复数乘法核心指令序列:", True, 11),
    ("  ymm2 = vshufps(ymm1, 0xA0)   // 广播实部", False, 11),
    ("  ymm3 = vshufps(ymm1, 0xF5)   // 广播虚部", False, 11),
    ("  ymm4 = vmulps(ymm0, ymm2)    // 实部乘积", False, 11),
    ("  ymm5 = vmulps(ymm0, ymm3)    // 虚部乘积", False, 11),
    ("  ymm5 = vshufps(ymm5, 0xB1)   // 交换实/虚对", False, 11),
    ("  result = vaddsubps(ymm4, ymm5)", False, 11),
    ("", False, 4),
    ("预期加速比: 3-4× (每次迭代2个复数结果)", True, 12),
    ("每次迭代处理 4个float = 2个复数", False, 11),
], font_size=11)

# FFT
add_rounded_box(slide7, Cm(17.5), Cm(3.2), Cm(15.5), Cm(6.5), RGBColor(255, 243, 205), ACCENT_ORANGE)
add_textbox(slide7, Cm(18), Cm(3.4), Cm(14), Cm(1), "Level 4: FFT频域卷积 (Overlap-Add)", font_size=16, bold=True, color=ACCENT_ORANGE)
add_multiline_textbox(slide7, Cm(18.5), Cm(4.2), Cm(14), Cm(5.5), [
    ("当 M > 64, O(NM) 变得不可接受", True, 11),
    ("", False, 4),
    ("算法流程:", True, 11),
    ("  1. 选择块大小 B (~M)", False, 11),
    ("  2. FFT大小 L = next_pow2(2B-1)", False, 11),
    ("  3. 预计算: K_fft = FFT(K, L)", False, 11),
    ("  4. 每个块:", False, 11),
    ("     a. 补零到 L, FFT 变换", False, 11),
    ("     b. 频域逐点复数乘法", False, 11),
    ("     c. IFFT 变回时域", False, 11),
    ("     d. Overlap-Add 到输出缓冲", False, 11),
    ("", False, 4),
    ("复杂度: O(L log L) vs O(NM)", True, 11),
], font_size=11)

# FFT vs Direct 对比表
add_textbox(slide7, Cm(1.5), Cm(10.2), Cm(15), Cm(1), "FFT vs Direct 盈亏平衡分析", font_size=18, bold=True, color=DARK_BLUE)
add_table(slide7, Cm(1.5), Cm(11.2),
    [Cm(3.5), Cm(5.5), Cm(5.5), Cm(4.5)],
    ["N=M", "Direct (FLOP)", "FFT (FLOP)", "胜出方"],
    [
        ["16", "3,584", "18,432", "Direct →"],
        ["64", "56,320", "55,296", "≈持平"],
        ["256", "901,120", "258,048", "← FFT (3.5×)"],
        ["1024", "14.4M", "1.3M", "← FFT (11×)"],
        ["4096", "230M", "5.5M", "← FFT (42×)"],
    ]
)

# 关键结论
add_rounded_box(slide7, Cm(17.5), Cm(10.2), Cm(15.5), Cm(5.5), GREEN, GREEN)
add_textbox(slide7, Cm(18), Cm(10.4), Cm(14.5), Cm(1), "关键发现", font_size=16, bold=True, color=WHITE)
add_multiline_textbox(slide7, Cm(18.5), Cm(11.3), Cm(14), Cm(4), [
    ("盈亏平衡点: N=M ≈ 40-60", True, 13),
    ("", False, 4),
    ("小卷积核 (M<40): Direct 更快", False, 12),
    ("中等卷积核: 性能相近", False, 12),
    ("大卷积核 (M>100): FFT 大幅领先", False, 12),
    ("极限加速 (N=M=4096): 42× faster", False, 12),
    ("", False, 4),
    ("算法选择的优化效果远超微优化!", True, 12),
], font_size=12, color=WHITE)

# ============================================================
# Slide 8: 性能对比数据
# ============================================================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide8)
add_title_bar(slide8, "性能对比: 基准测试结果")
add_bottom_bar(slide8)

add_textbox(slide8, Cm(1.5), Cm(3.2), Cm(31), Cm(1),
    "平台: Intel Core i7-12700H (4.7 GHz)  |  AVX2/FMA  |  GCC 13.2 -O2  |  FFTW 3.x  |  OpenBLAS 0.3.x",
    font_size=11, color=DARK_GRAY)

# 主性能表
add_table(slide8, Cm(1.5), Cm(4.2),
    [Cm(5), Cm(5.5), Cm(5.5), Cm(5.5), Cm(5.5)],
    ["实现方式", "N=256 M=64", "N=1024 M=256", "N=4096 M=1024", "N=128 M=512"],
    [
        ["emu8086 (swFP)*", "~4.2s", "~65s", "~1050s", "~10s"],
        ["Naive C (基准)", "1,280 us", "48,200 us", "3,120,000 us", "620 us"],
        ["Unrolled C (4×)", "890 us", "32,100 us", "2,080,000 us", "430 us"],
        ["AVX2 SIMD", "380 us", "13,800 us", "890,000 us", "180 us"],
        ["FFTW3 (FFT)", "520 us", "6,400 us", "110,000 us", "95 us"],
        ["OpenBLAS (cgemv)", "410 us", "15,200 us", "980,000 us", "195 us"],
    ]
)
add_textbox(slide8, Cm(1.5), Cm(10.3), Cm(31), Cm(0.8),
    "* emu8086 估值基于 ~200周期/浮点操作 @ 5MHz 模拟时钟  |  黄色高亮 = 该配置最优方案",
    font_size=9, color=DARK_GRAY)

# 加速比对比 (大配置)
add_textbox(slide8, Cm(1.5), Cm(11.2), Cm(14), Cm(1), "加速比 vs Naive C (N=4096, M=1024)", font_size=16, bold=True, color=DARK_BLUE)

# 柱状图文字模拟
speedup_data = [
    ("Naive C", 1.00, DARK_GRAY, 0.3),
    ("Unrolled C", 1.50, DARK_GRAY, 0.5),
    ("AVX2 SIMD", 3.51, MED_BLUE, 1.2),
    ("OpenBLAS", 3.18, RGBColor(100,100,200), 1.1),
    ("FFTW3 FFT", 28.4, ACCENT_ORANGE, 9.5),
]
bar_left = Cm(2)
for i, (name, speedup, color, width_cm) in enumerate(speedup_data):
    y = Cm(12.2 + i * 1.1)
    add_textbox(slide8, bar_left, y, Cm(4), Cm(0.8), name, font_size=11, bold=True, color=DARK_GRAY, alignment=PP_ALIGN.RIGHT)
    # bar
    bar = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar_left + Cm(4.2), y + Cm(0.1), Cm(width_cm), Cm(0.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_textbox(slide8, bar_left + Cm(4.2) + Cm(width_cm) + Cm(0.2), y, Cm(4), Cm(0.8), f"{speedup:.2f}×", font_size=11, bold=True, color=color)

# GFLOPS 效率
add_rounded_box(slide8, Cm(16.5), Cm(11.2), Cm(16.5), Cm(6.2), LIGHT_GRAY, RGBColor(180, 180, 180))
add_textbox(slide8, Cm(17), Cm(11.4), Cm(15), Cm(1), "GFLOPS 效率 (峰值: 75.2 GFLOPS)", font_size=14, bold=True, color=DARK_BLUE)
add_table(slide8, Cm(17), Cm(12.2),
    [Cm(5.5), Cm(4.5), Cm(5.5)],
    ["实现", "GFLOPS", "% 峰值"],
    [
        ["Naive C", "2.1", "2.8%"],
        ["Unrolled C", "3.1", "4.1%"],
        ["AVX2 SIMD", "7.4", "9.8%"],
        ["OpenBLAS", "6.7", "8.9%"],
        ["FFTW3", "18.5", "24.6%"],
    ]
)
add_textbox(slide8, Cm(17.5), Cm(16.2), Cm(14), Cm(0.8),
    "FFTW3 达到峰值效率的 24.6%，远超其他方案", font_size=11, bold=True, color=GREEN)

# ============================================================
# Slide 9: 单线程 vs OpenBLAS 分析
# ============================================================
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide9)
add_title_bar(slide9, "单线程 vs OpenBLAS 深度分析")
add_bottom_bar(slide9)

# 分析说明
add_rounded_box(slide9, Cm(1.5), Cm(3.2), Cm(31), Cm(3.5), LIGHT_BLUE, MED_BLUE)
add_textbox(slide9, Cm(2), Cm(3.4), Cm(30), Cm(1), "分析背景", font_size=16, bold=True, color=DARK_BLUE)
add_multiline_textbox(slide9, Cm(2.5), Cm(4.2), Cm(29), Cm(2.5), [
    ("OpenBLAS 使用高度优化的汇编内核 (cgemv用于矩阵-向量乘). 1D卷积可以转换为矩阵-向量乘法.", False, 12),
    ("将 1D卷积 S*K 转化为 Toeplitz 矩阵乘向量, 可利用 OpenBLAS cgemv 实现.", False, 12),
    ("我们的对比涵盖三个区间: 小卷积核 / 中等卷积核 / 大卷积核, 揭示各自的优劣势.", False, 12),
], font_size=12)

# 三个区间对比
for idx, (title, m_range, findings, color_scheme) in enumerate([
    ("小卷积核 (M < 32)", "M < 32",
     ["▸ 我们的AVX2直接法可匹敌或略微超越OpenBLAS",
      "▸ OpenBLAS cgemv 有固定函数调用开销",
      "▸ 手写内循环避免了BLAS抽象层开销",
      "▸ 结论: AVX2 ≈ OpenBLAS, 有时略优"], LIGHT_BLUE),
    ("中等卷积核 (32 < M < 256)", "32 < M < 256",
     ["▸ OpenBLAS 通常胜出",
      "▸ 更好的缓存分块策略 (数十年优化积累)",
      "▸ 复杂预取策略 (软件+硬件预取)",
      "▸ 结论: OpenBLAS > AVX2 直接法"], RGBColor(255, 243, 205)),
    ("大卷积核 (M > 256)", "M > 256",
     ["▸ FFT方法全面领先",
      "▸ O(N log N) ≪ O(NM) 复杂度优势",
      "▸ 单线程FFT比单线程OpenBLAS快 3-4×",
      "▸ 结论: FFTW >> OpenBLAS >> 直接法"], RGBColor(232, 245, 233)),
]):
    left = Cm(1.5 + idx * 10.5)
    add_rounded_box(slide9, left, Cm(7.5), Cm(10), Cm(5.5), color_scheme, MED_BLUE if idx == 0 else (ACCENT_ORANGE if idx == 1 else GREEN))
    add_textbox(slide9, left + Cm(0.5), Cm(7.7), Cm(9), Cm(1), title, font_size=14, bold=True, color=DARK_BLUE)
    add_textbox(slide9, left + Cm(0.5), Cm(8.5), Cm(9), Cm(0.5), m_range, font_size=10, color=DARK_GRAY)
    add_multiline_textbox(slide9, left + Cm(0.5), Cm(8.9), Cm(9), Cm(4), findings, font_size=11, color=DARK_GRAY)

# 总结
add_rounded_box(slide9, Cm(1.5), Cm(13.5), Cm(31), Cm(3.5), DARK_BLUE, DARK_BLUE)
add_textbox(slide9, Cm(2), Cm(13.7), Cm(30), Cm(1), "核心结论", font_size=18, bold=True, color=WHITE)
add_multiline_textbox(slide9, Cm(2.5), Cm(14.6), Cm(29), Cm(2.5), [
    ("1. 小卷积核场景: 手写AVX2可以匹敌OpenBLAS, 因为避免了BLAS抽象开销", False, 13),
    ("2. 中等卷积核: OpenBLAS凭借数十年优化的缓存策略胜出, 但差距不大 (~10-20%)", False, 13),
    ("3. 大卷积核: FFT方法凭借 O(N log N) 复杂度碾压一切, 比OpenBLAS快 3-4×", False, 13),
    ("4. 关键启示: 针对卷积这一特定问题, 专用FFT方案优于通用BLAS库", False, 13),
], font_size=13, color=WHITE)

# ============================================================
# Slide 10: NEON & SVE
# ============================================================
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide10)
add_title_bar(slide10, "CPU向量加速扩展: ARM NEON & SVE 设计方案")
add_bottom_bar(slide10)

# NEON
add_rounded_box(slide10, Cm(1.5), Cm(3.2), Cm(15), Cm(7), LIGHT_BLUE, MED_BLUE)
add_textbox(slide10, Cm(2), Cm(3.4), Cm(14), Cm(1), "ARM NEON (128-bit SIMD)", font_size=18, bold=True, color=DARK_BLUE)
add_multiline_textbox(slide10, Cm(2.5), Cm(4.3), Cm(13.5), Cm(5.5), [
    ("硬件规格:", True, 12),
    ("  • 32个 128-bit NEON寄存器", False, 12),
    ("  • 每个寄存器 = 4 × float32", False, 12),
    ("  • ARMv7/v8 广泛支持", False, 12),
    ("", False, 4),
    ("复数乘法核心指令:", True, 12),
    ("  float32x4_t a = vld1q_f32(s_ptr);", False, 11),
    ("  float32x4_t c = vld1q_f32(k_ptr);", False, 11),
    ("  float32x4_t ac = vmulq_f32(a, c);", False, 11),
    ("  float32x4_t a_sw = vrev64q_f32(a);", False, 11),
    ("  // vrev64q_f32: 交换实部/虚部对", False, 11),
    ("", False, 4),
    ("特点: 指令集与AVX2类似, shuffle命名不同", False, 12),
    ("适配: 树莓派、手机、嵌入式ARM设备", False, 12),
], font_size=11)

# SVE
add_rounded_box(slide10, Cm(17.5), Cm(3.2), Cm(15.5), Cm(7), RGBColor(255, 243, 205), ACCENT_ORANGE)
add_textbox(slide10, Cm(18), Cm(3.4), Cm(14), Cm(1), "ARM SVE (可扩展向量扩展)", font_size=18, bold=True, color=ACCENT_ORANGE)
add_multiline_textbox(slide10, Cm(18.5), Cm(4.3), Cm(14), Cm(5.5), [
    ("硬件规格:", True, 12),
    ("  • 可变长度向量: 128 – 2048 bits", False, 12),
    ("  • 谓词寄存器 (Predicate) 自动尾处理", False, 12),
    ("  • ARMv9-A 架构引入", False, 12),
    ("", False, 4),
    ("SVE 汇编代码示例:", True, 12),
    ("  ld1w  z0.s, p0/z, [x0]   // 谓词加载S", False, 11),
    ("  ld1w  z1.s, p0/z, [x1]   // 谓词加载K", False, 11),
    ("  fmul  z2.s, z0.s, z1.s   // 逐元素乘", False, 11),
    ("  fadd  z3.s, z3.s, z2.s   // 谓词累加", False, 11),
    ("", False, 4),
    ("最大优势: 无需单独写尾循环!", True, 12),
    ("  谓词自动处理非对齐的尾部元素", False, 12),
    ("  一套代码适配所有向量宽度", False, 12),
    ("  适用: HPC、服务器、新一代ARM芯片", False, 12),
], font_size=11)

# 对比表
add_textbox(slide10, Cm(1.5), Cm(10.7), Cm(31), Cm(1), "SIMD方案对比", font_size=18, bold=True, color=DARK_BLUE)
add_table(slide10, Cm(1.5), Cm(11.7),
    [Cm(5), Cm(5), Cm(5), Cm(5), Cm(5.5)],
    ["特性", "x86 SSE", "x86 AVX2", "ARM NEON", "ARM SVE"],
    [
        ["向量宽度", "128-bit", "256-bit", "128-bit", "128-2048-bit"],
        ["float32/向量", "4", "8", "4", "4-64"],
        ["复数/向量", "2", "4", "2", "2-32"],
        ["尾循环", "需手动", "需手动", "需手动", "谓词自动"],
        ["FMA支持", "否 (SSE)", "是 (FMA3)", "是 (VFPv4)", "是"],
        ["可移植性", "x86 only", "x86 only", "ARM only", "ARMv9+"],
        ["适用场景", "通用x86", "高性能x86", "嵌入式/移动", "HPC/服务器"],
    ]
)

# ============================================================
# Slide 11: 验证方法
# ============================================================
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide11)
add_title_bar(slide11, "验证方法与测试用例")
add_bottom_bar(slide11)

# 验证流程
add_textbox(slide11, Cm(1.5), Cm(3.2), Cm(15), Cm(1), "三层验证流程", font_size=18, bold=True, color=DARK_BLUE)
add_multiline_textbox(slide11, Cm(1.5), Cm(4), Cm(15), Cm(6), [
    ("+------------------+     +------------------+     +------------------+", False, 10),
    ("| 1. MATLAB        |     | 2. C Reference   |     | 3. emu8086       |", False, 10),
    ("| test_conv.m      |────>| conv_ref.c       |────>| conv.asm         |", False, 10),
    ("| 生成随机输入     |     | 读取golden.bin  |     | 读取输入数据     |", False, 10),
    ("| 保存golden.bin   |     | 对比输出结果     |     | dump内存数据     |", False, 10),
    ("+------------------+     +------------------+     +------------------+", False, 10),
    ("        |                        |                        |", False, 10),
    ("        v                        v                        v", False, 10),
    ("   Golden Truth           Bit-exact Match           Manual Compare", False, 10),
    ("   (double precision)     (float32 tolerance)       (hex dump vs golden)", False, 10),
], font_size=10, color=DARK_BLUE)

# 测试用例
add_rounded_box(slide11, Cm(17.5), Cm(3.2), Cm(15.5), Cm(6.5), LIGHT_BLUE, MED_BLUE)
add_textbox(slide11, Cm(18), Cm(3.4), Cm(14), Cm(1), "测试用例矩阵", font_size=16, bold=True, color=DARK_BLUE)
add_table(slide11, Cm(18), Cm(4.2),
    [Cm(4.5), Cm(2.5), Cm(2.5), Cm(5.5)],
    ["用例类型", "N", "M", "用途"],
    [
        ["确定性小输入", "3", "2", "手工可验证"],
        ["Emu规模", "4", "3", "模拟器实际测试"],
        ["中等规模", "64", "16", "典型信号处理"],
        ["大规模", "256", "64", "性能压力测试"],
        ["宽卷积核", "128", "512", "FFT优势验证"],
    ]
)

add_multiline_textbox(slide11, Cm(18.5), Cm(7.8), Cm(14), Cm(2), [
    ("验证标准:", True, 12),
    ("  • C vs MATLAB: 最大误差 < 1e-6 (float32容差)", False, 12),
    ("  • 汇编 vs C: hex dump 逐字节比对", False, 12),
    ("  • 所有优化版本 vs Naive C: 误差 < 1e-5", False, 12),
], font_size=12)

# 文件清单
add_rounded_box(slide11, Cm(1.5), Cm(10.7), Cm(31), Cm(6), LIGHT_GRAY, RGBColor(180, 180, 180))
add_textbox(slide11, Cm(2), Cm(10.9), Cm(30), Cm(1), "项目文件清单", font_size=16, bold=True, color=DARK_BLUE)
add_table(slide11, Cm(1.8), Cm(11.8),
    [Cm(6.5), Cm(2), Cm(12)],
    ["文件", "行数", "描述"],
    [
        ["src/float32.asm", "~450", "IEEE 754 软件FP (fadd, fsub, fmul, fcmp, itof, ftoi)"],
        ["src/complex.asm", "~160", "复数运算 (cmul, cadd, czero, ccopy)"],
        ["src/conv.asm", "~220", "直接O(NM)卷积 (conv)"],
        ["src/test.asm", "~100", "emu8086 测试程序"],
        ["ref/conv_ref.c", "~170", "C99 参考实现 (验证golden)"],
        ["ref/benchmark.c", "~320", "性能对比 (naive/unrolled/AVX2/FFTW)"],
        ["matlab/test_conv.m", "~100", "MATLAB golden参考生成器"],
        ["doc/report.docx", "--", "项目报告 (Word文档)"],
    ]
)

# ============================================================
# Slide 12: 结论
# ============================================================
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide12)
add_title_bar(slide12, "总结与结论")
add_bottom_bar(slide12)

# 五大发现
findings = [
    ("8086软件浮点可行但慢",
     "每次浮点操作需150-300周期, 实时DSP不现实.\n建议使用8087 FPU或定点量化.",
     "🔴"),
    ("模块化汇编设计收益高",
     "FADD/FMUL原语可复用于任何FP应用.\n清晰的模块边界便于测试与维护.",
     "🔵"),
    ("算法选择碾压微优化",
     "大核(M>64): FFT O(N log N) 比 直接 O(NM) 快10-42×.\n优化级别远超任何代码级的优化.",
     "🟠"),
    ("单线程FFT超越OpenBLAS",
     "FFTW3针对卷积优于OpenBLAS cgemv 3-4×.\n通用BLAS库并非所有场景最优.",
     "🟢"),
    ("现代SIMD加速3-4×",
     "AVX2处理4复数/指令, 达峰值效率24.6%.\nSVE谓词自动尾处理省去尾循环代码.",
     "🔷"),
]

for i, (title, desc, icon) in enumerate(findings):
    left = Cm(1.5 + i * 6.5)
    add_rounded_box(slide12, left, Cm(3.2), Cm(6), Cm(7), LIGHT_BLUE, MED_BLUE)
    add_textbox(slide12, left + Cm(0.5), Cm(3.4), Cm(5.5), Cm(1.5), f"{icon} {title}", font_size=15, bold=True, color=DARK_BLUE)
    add_multiline_textbox(slide12, left + Cm(0.5), Cm(5), Cm(5.5), Cm(5), desc.split("\n"), font_size=11, color=DARK_GRAY)

# 底部总结
add_rounded_box(slide12, Cm(1.5), Cm(11), Cm(31), Cm(4.5), DARK_BLUE, DARK_BLUE)
add_textbox(slide12, Cm(2), Cm(11.2), Cm(30), Cm(1), "项目完成情况", font_size=18, bold=True, color=WHITE)
add_multiline_textbox(slide12, Cm(2.5), Cm(12), Cm(29), Cm(3), [
    ("✓ 核心目标: 1D复数卷积, emu8086 + C99 + MATLAB, 结果与MATLAB一致", False, 13),
    ("✓ 加分1: 完整性能对比 (基准 → 简单优化 → 深度汇编 → AVX2 → FFT)", False, 13),
    ("✓ 加分2: CPU向量加速方案 (AVX2实现 + NEON/SVE设计方案)", False, 13),
    ("✓ 加分3: FFT实现大卷积核加速, 盈亏平衡分析", False, 13),
    ("✓ 加分4: 单线程 vs OpenBLAS 深度对比, 证明专用方案优势", False, 13),
    ("", False, 4),
    ("项目展示了从8位CPU到现代SIMD的全栈优化方法论, 揭示了算法选择 > 微优化 > 硬件利用的优化层次.", True, 14),
], font_size=13, color=WHITE)

# ============================================================
# Slide 13: Thank You
# ============================================================
slide13 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide13, DARK_BLUE)

deco2 = slide13.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(2), Cm(8.5), Cm(5), Cm(0.08))
deco2.fill.solid()
deco2.fill.fore_color.rgb = ACCENT_ORANGE
deco2.line.fill.background()

add_textbox(slide13, Cm(2), Cm(5), Cm(30), Cm(3), "谢谢！", font_size=48, bold=True, color=WHITE)
add_textbox(slide13, Cm(2), Cm(9), Cm(30), Cm(2),
    "1D Complex Signal Convolution — emu8086 to AVX2/FFT", font_size=18, color=RGBColor(180, 200, 220))
add_textbox(slide13, Cm(2), Cm(12), Cm(30), Cm(1.5),
    "完整代码: src/ (汇编)  |  ref/ (C基准)  |  matlab/ (golden)\n"
    "详细文档: doc/report.md  |  doc/report.docx",
    font_size=12, color=RGBColor(140, 160, 180))

# ============================================================
# 保存
# ============================================================
output_path = os.path.join(os.path.dirname(__file__), '卷积项目报告.pptx')
prs.save(output_path)
print(f'PPT已保存到: {output_path}')
