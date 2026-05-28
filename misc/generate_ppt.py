from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BASE = r"C:\Users\24600\Desktop\大二下\计组\todo"

# ============================================================
# SLIDE 1 - 封面
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# Background
bg = slide1.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
bg.line.fill.background()

# Title
tb = slide1.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.3), Inches(1.5))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "emu8086 汇编语言课程设计"
p.alignment = PP_ALIGN.CENTER
run = p.runs[0]
run.font.size = Pt(48)
run.font.color.rgb = RGBColor(0xE9, 0x45, 0x60)
run.font.bold = True

# Subtitle
tb2 = slide1.shapes.add_textbox(Inches(1), Inches(3.3), Inches(11.3), Inches(2.0))
tf2 = tb2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.text = "成绩判定器  &  矩阵乘法"
p2.alignment = PP_ALIGN.CENTER
run2 = p2.runs[0]
run2.font.size = Pt(40)
run2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
run2.font.bold = True

# Line
line = slide1.shapes.add_shape(1, Inches(3), Inches(5.0), Inches(7.3), Inches(0.03))
line.fill.solid()
line.fill.fore_color.rgb = RGBColor(0xE9, 0x45, 0x60)
line.line.fill.background()

# Info
tb3 = slide1.shapes.add_textbox(Inches(2), Inches(5.5), Inches(9.3), Inches(1.2))
tf3 = tb3.text_frame
tf3.word_wrap = True
p3 = tf3.paragraphs[0]
p3.text = "基于 Intel 8086 指令集  |  整数定点运算  |  浮点数处理  |  矩阵运算"
p3.alignment = PP_ALIGN.CENTER
run3 = p3.runs[0]
run3.font.size = Pt(20)
run3.font.color.rgb = RGBColor(0xAA, 0xAA, 0xBB)

# ============================================================
# SLIDE 2 - 成绩判定器
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])

# Title bar
bar2 = slide2.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.1))
bar2.fill.solid()
bar2.fill.fore_color.rgb = RGBColor(0xE9, 0x45, 0x60)
bar2.line.fill.background()
tb_title2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.8))
tf_title2 = tb_title2.text_frame
p_title2 = tf_title2.paragraphs[0]
p_title2.text = "任务一：成绩判定器 (Grade Determiner)"
p_title2.alignment = PP_ALIGN.LEFT
run_title2 = p_title2.runs[0]
run_title2.font.size = Pt(36)
run_title2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
run_title2.font.bold = True

# Left panel - 原理
left_box = slide2.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(5.8), Inches(5.5))
lf = left_box.text_frame
lf.word_wrap = True

def add_para(tf, text, size=18, bold=False, color=RGBColor(0x33,0x33,0x33), space_after=Pt(8)):
    p = tf.add_paragraph()
    p.text = text
    p.space_after = space_after
    p.alignment = PP_ALIGN.LEFT
    if text:
        run = p.runs[0]
        run.font.size = Pt(size) if isinstance(size, int) else size
        run.font.bold = bold
        run.font.color.rgb = color
    return p

# First paragraph (already exists)
p_first = lf.paragraphs[0]
p_first.text = "核心功能"
p_first.alignment = PP_ALIGN.LEFT
run = p_first.runs[0]
run.font.size = Pt(26)
run.font.bold = True
run.font.color.rgb = RGBColor(0xE9, 0x45, 0x60)

add_para(lf, "")  # spacer
add_para(lf, "将 0-100 分数（支持小数如 89.5）转换为等级 A~D：")
add_para(lf, "  A: Score >= 90          优秀")
add_para(lf, "  B: 75 <= Score < 90     良好")
add_para(lf, "  C: 60 <= Score < 75     及格")
add_para(lf, "  D: Score < 60            不及格")
add_para(lf, "")  # spacer

p_key = add_para(lf, "关键技术点", size=22, bold=True, color=RGBColor(0xE9,0x45,0x60))
p_key = add_para(lf, "")
add_para(lf, "1. 字符串解析：支持整数与一位小数输入（如 89.5）")
add_para(lf, "2. 定点数处理：分数 * 10 转为整数存储，避免浮点运算")
add_para(lf, "3. 逐字符 DFA 解析：整数部分 → '.' → 小数部分")
add_para(lf, "4. 范围校验：>100、<0、非法字符均报错并提示重试")
add_para(lf, "5. 循环执行：支持 Y/N 再次输入，无需重启程序")

# Right panel - image
img_path1 = os.path.join(BASE, "成绩判定器", "成绩判定器.png")
if os.path.exists(img_path1):
    pic = slide2.shapes.add_picture(img_path1, Inches(6.8), Inches(1.4), Inches(6.0), Inches(5.0))

# ============================================================
# SLIDE 3 - 矩阵乘法
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])

# Title bar
bar3 = slide3.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.1))
bar3.fill.solid()
bar3.fill.fore_color.rgb = RGBColor(0x16, 0x6D, 0x8C)
bar3.line.fill.background()
tb_title3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.8))
tf_title3 = tb_title3.text_frame
p_title3 = tf_title3.paragraphs[0]
p_title3.text = "任务二：矩阵乘法 (Matrix Multiplication)"
p_title3.alignment = PP_ALIGN.LEFT
run_title3 = p_title3.runs[0]
run_title3.font.size = Pt(36)
run_title3.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
run_title3.font.bold = True

# Left panel
left3 = slide3.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(5.8), Inches(5.8))
lf3 = left3.text_frame
lf3.word_wrap = True

p3_first = lf3.paragraphs[0]
p3_first.text = "核心功能"
p3_first.alignment = PP_ALIGN.LEFT
run3f = p3_first.runs[0]
run3f.font.size = Pt(26)
run3f.font.bold = True
run3f.font.color.rgb = RGBColor(0x16, 0x6D, 0x8C)

add_para(lf3, "")
add_para(lf3, "实现任意维度矩阵乘法 C[MxO] = A[MxN] × B[NxO]")
add_para(lf3, "")
add_para(lf3, "Demo: A[2x3] × B[3x2] = C[2x2]")
add_para(lf3, "")
add_para(lf3, "     | 1  2  3 |          | 1  2 |")
add_para(lf3, "A =  | 4  5  6 |    B =   | 3  4 |")
add_para(lf3, "                         | 5  6 |")
add_para(lf3, "")
add_para(lf3, "        | 1*1+2*3+3*5   1*2+2*4+3*6 |   | 22  28 |")
add_para(lf3, "C = A×B=| 4*1+5*3+6*5   4*2+5*4+6*6 | = | 49  64 |")
add_para(lf3, "")

add_para(lf3, "关键技术点", size=22, bold=True, color=RGBColor(0x16,0x6D,0x8C))
add_para(lf3, "")
add_para(lf3, "1. 三层循环：r(行) × c(列) × k(内积) 经典矩阵乘法")
add_para(lf3, "2. 定点数：所有元素 ×100 存储，避免 8087 协处理器")
add_para(lf3, "3. 整数 MUL/DIV：乘积 /100 后累加，保证精度")
add_para(lf3, "4. 通用显示：定点数换行打印，支持符号和两位小数")

# Right side images
img_path2 = os.path.join(BASE, "矩阵乘法", "矩阵乘法原理.png")
if os.path.exists(img_path2):
    pic2 = slide3.shapes.add_picture(img_path2, Inches(6.6), Inches(1.3), Inches(6.3), Inches(2.6))

img_path3 = os.path.join(BASE, "矩阵乘法", "输出结果.png")
if os.path.exists(img_path3):
    pic3 = slide3.shapes.add_picture(img_path3, Inches(6.6), Inches(4.1), Inches(6.3), Inches(3.0))

# ============================================================
# SLIDE 4 - 总结
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])

# Title bar
bar4 = slide4.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.1))
bar4.fill.solid()
bar4.fill.fore_color.rgb = RGBColor(0x2D, 0x2D, 0x4A)
bar4.line.fill.background()
tb_title4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.8))
tf_title4 = tb_title4.text_frame
p_title4 = tf_title4.paragraphs[0]
p_title4.text = "总结与收获"
p_title4.alignment = PP_ALIGN.LEFT
run_title4 = p_title4.runs[0]
run_title4.font.size = Pt(36)
run_title4.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
run_title4.font.bold = True

# Two columns
# Left column
left4 = slide4.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.5))
lf4 = left4.text_frame
lf4.word_wrap = True

p4f = lf4.paragraphs[0]
p4f.text = "成绩判定器"
p4f.alignment = PP_ALIGN.LEFT
r4f = p4f.runs[0]
r4f.font.size = Pt(24)
r4f.font.bold = True
r4f.font.color.rgb = RGBColor(0xE9, 0x45, 0x60)

add_para(lf4, "")
add_para(lf4, "  DOS中断 INT 21h 实现输入输出")
add_para(lf4, "  缓冲输入 (AH=0Ah) 读取字符串")
add_para(lf4, "  逐字符解析实现 DFA 状态机")
add_para(lf4, "  分支跳转 (cmp/jae/je) 判断等级")
add_para(lf4, "  循环控制实现多次输入")

add_para(lf4, "")
add_para(lf4, "汇编核心能力", size=22, bold=True, color=RGBColor(0x16,0x6D,0x8C))
add_para(lf4, "")
add_para(lf4, "  寄存器操作与寻址方式")
add_para(lf4, "  MUL/DIV 定点整数运算代替浮点")
add_para(lf4, "  子程序调用与栈帧管理")
add_para(lf4, "  字符串格式化显示")

# Right column
right4 = slide4.shapes.add_textbox(Inches(7.0), Inches(1.5), Inches(5.5), Inches(5.5))
rf4 = right4.text_frame
rf4.word_wrap = True

pr4f = rf4.paragraphs[0]
pr4f.text = "矩阵乘法"
pr4f.alignment = PP_ALIGN.LEFT
rr4f = pr4f.runs[0]
rr4f.font.size = Pt(24)
rr4f.font.bold = True
rr4f.font.color.rgb = RGBColor(0x16, 0x6D, 0x8C)

add_para(rf4, "")
add_para(rf4, "  三层嵌套循环实现矩阵乘法")
add_para(rf4, "  一维数组模拟二维矩阵(行优先)")
add_para(rf4, "  index = (r*cols + c) * sizeof(word)")
add_para(rf4, "  定点数 100 倍缩放，MUL+DIV 还原")
add_para(rf4, "  通用矩阵显示子程序")

add_para(rf4, "")
add_para(rf4, "emu8086 工具链", size=22, bold=True, color=RGBColor(0xE9,0x45,0x60))
add_para(rf4, "")
add_para(rf4, "  emu8086 模拟器 + 内置调试器")
add_para(rf4, "  寄存器/内存窗口观察运行状态")
add_para(rf4, "  单步执行验证循环与分支逻辑")
add_para(rf4, "  .com / .exe 两种输出格式")

# Output
out_path = os.path.join(BASE, "计组课程设计汇报.pptx")
prs.save(out_path)
print(f"PPT saved to: {out_path}")
